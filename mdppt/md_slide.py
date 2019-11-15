import os

from pptx.util import Inches
from pygments import highlight
from pygments.formatters.img import ImageFormatter
from pygments.lexers.python import PythonLexer

from mdppt.md_dot import MDDot
from mdppt.tw_layout import TWLayout

__location__ = os.path.realpath(
    os.path.join(os.getcwd(), os.path.dirname(__file__)))


class MdSlider:
    def __init__(self, prs) -> None:
        self.code_count = 0
        self.dot_image_count = 0
        self.slides = prs.slides

        self.SLIDE_WIDTH = prs.slide_height
        self.SLIDE_HEIGHT = prs.slide_width
        self.current_slide = None

        self.tw_layout = TWLayout(prs.slide_masters)
        super().__init__()

    def add_normal_slide(self, content):
        black_layout = self.tw_layout.get_layout_by_level(content.header_level)
        normal_slide = self.slides.add_slide(black_layout)
        title = normal_slide.placeholders[0]
        title.text = content.header

        content_section = normal_slide.placeholders[1]
        content_section.text = content.paragraph

        self.current_slide = normal_slide
        return self.current_slide

    def add_list_slide(self, header, list):
        layout = self.tw_layout.get_list_layout()
        list_slide = self.slides.add_slide(layout)

        title = list_slide.placeholders[0]
        title.text = header

        content_section = list_slide.placeholders[1]

        for item in list:
            content_section.text = content_section.text + item + "\n"

        return layout

    def add_quote_slide(self, text):
        layout = self.tw_layout.get_quote_layout()
        quote_slide = self.slides.add_slide(layout)

        quote_section = quote_slide.placeholders[0]
        quote_section.text = text

        self.current_slide = quote_slide
        return self.current_slide

    def add_code(self, code, position):
        if position == "right":
            values = self.build_code_image(code)
            self.write_image_to_file(values, self.code_count)
            self.insert_image_by_type(self.code_count)

    def add_image(self, src, position):
        self.current_slide.placeholders[1].text = ""
        if position == "right":
            top = Inches(1.5)
            left = self.SLIDE_WIDTH / 2
            self.current_slide.shapes.add_picture(src, left, top, left)

    def add_dot(self, dot, position):
        md_dot = MDDot()
        md_dot.build_graph(dot, self.dot_image_count)
        self.insert_image_by_type(self.code_count, image_type="dot")

    def insert_image_by_type(self, code_count, image_type='code'):
        top = Inches(1.5)
        left = self.SLIDE_WIDTH / 2
        imagePath = os.path.join(__location__, 'images/' + image_type + code_count.__str__() + '.png')
        self.current_slide.shapes.add_picture(imagePath, left, top, left)

    def build_code_image(self, code):
        return highlight(code, PythonLexer(), ImageFormatter(noclasses=True,
                                                             image_format='bmp',
                                                             image_pad='32',
                                                             line_number_bg='#fff',
                                                             line_numbers=False,
                                                             line_number_fg='#fff',
                                                             font_size='64',
                                                             font_name='Inconsolata Bold for Powerline'))

    def write_image_to_file(self, values, code=None):
        image_dir = os.path.join(__location__, 'images')
        if not os.path.exists(image_dir):
            os.mkdir(image_dir, 0o755)

        fo = open(os.path.join(__location__, 'images/code' + code.__str__() + '.png'), 'wb')
        fo.write(values)
        fo.close()

    def add_table(self, table):
        table_height = len(table)
        table_width = len(table[0])

        x, y, cx, cy = Inches(2), Inches(2), Inches(4), Inches(1.5)
        shape = self.current_slide.shapes.add_table(table_height, table_width, x, y, cx, cy)
        for r_index, row in enumerate(table):
            for c_index, cell in enumerate(row):
                shape.table.cell(r_index, c_index).text = cell
