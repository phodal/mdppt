from pptx import Presentation
from pptx.util import Pt

def runMdppt():
    f = open(r'templates/tw-slides2.pptx', 'rb')
    prs = Presentation(f)
    f.close()

    slide = prs.slides.add_slide(prs.slide_masters[0].slide_layouts[0])

    title = slide.placeholders[0]
    title.text = "REPORT"

    prs.save('test.pptx')

if __name__ == "__main__":
    runMdppt()
